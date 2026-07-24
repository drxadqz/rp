from __future__ import annotations

import json
import sys
from pathlib import Path

sys.path.insert(0, r"E:\perception\friction_affordance_field\tmp\pydeps_ppt")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path.home() / "Desktop" / "RSCD视觉路面摩擦可供性估计_最近进展_完整流程公式解释版.pptx"
ASSET_DIR = Path(r"E:\perception\friction_affordance_field\presentation_build\recent_progress_formula_assets")
METRICS = Path(
    r"E:\perception_outputs\rscd_surface_classification"
    r"\c3_farnet_official_anchor_source_reliable_router_s7_fulltest_20260708\fast_test\metrics.json"
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 250, 252)
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
TEAL = RGBColor(13, 148, 136)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)
AMBER = RGBColor(217, 119, 6)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(203, 213, 225)
SOFT = RGBColor(241, 245, 249)


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_run_font(run, size=15, bold=False, color=SLATE, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def fill_text(shape, text, size=15, color=SLATE, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.08):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for pi, raw in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = raw
        set_run_font(run, size=size, bold=bold, color=color)


def add_text(slide, x, y, w, h, text, size=15, color=SLATE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    fill_text(shape, text, size=size, color=color, bold=bold, align=align)
    return shape


def add_title(slide, title, subtitle=None, tag=None):
    add_text(slide, 0.55, 0.23, 9.6, 0.48, title, size=25, color=NAVY, bold=True)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.80), Inches(1.30), Inches(0.05)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = TEAL
    slide.shapes[-1].line.fill.background()
    if subtitle:
        add_text(slide, 0.55, 0.90, 10.3, 0.32, subtitle, size=12.5, color=MUTED)
    if tag:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.95), Inches(0.32), Inches(1.82), Inches(0.34))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb("e0f2fe")
        box.line.color.rgb = rgb("bae6fd")
        fill_text(box, tag, size=10.5, color=rgb("075985"), bold=True, align=PP_ALIGN.CENTER)
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_panel(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    return shape


def add_panel_text(slide, x, y, w, h, header, body, header_color=TEAL, body_size=13.3):
    add_panel(slide, x, y, w, h)
    add_text(slide, x + 0.12, y + 0.10, w - 0.24, 0.25, header, size=12.5, color=header_color, bold=True)
    add_text(slide, x + 0.12, y + 0.40, w - 0.24, h - 0.50, body, size=body_size, color=SLATE)


def add_metric(slide, x, y, w, title, value, color):
    add_panel(slide, x, y, w, 0.76)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(0.76))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, x + 0.22, y + 0.12, w - 0.30, 0.26, value, size=19, color=color, bold=True)
    add_text(slide, x + 0.22, y + 0.43, w - 0.30, 0.20, title, size=9.5, color=SLATE)


def add_bullet_panel(slide, x, y, w, h, header, bullets, color=TEAL, size=13.2):
    add_panel(slide, x, y, w, h)
    add_text(slide, x + 0.14, y + 0.10, w - 0.28, 0.28, header, size=13.2, color=color, bold=True)
    text = "\n".join(f"• {b}" for b in bullets)
    add_text(slide, x + 0.17, y + 0.45, w - 0.30, h - 0.54, text, size=size, color=SLATE)


def render_formula(latex: str, name: str, fontsize: int = 19, width: float = 8.8, height: float = 0.58) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / f"{name}.png"
    if path.exists():
        return path
    plt.rcParams["mathtext.fontset"] = "stix"
    plt.rcParams["font.family"] = "DejaVu Serif"
    fig = plt.figure(figsize=(width, height), dpi=230)
    fig.patch.set_alpha(0)
    fig.text(0.5, 0.50, f"${latex}$", ha="center", va="center", fontsize=fontsize, color="#0f172a")
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    return path


def add_formula(slide, latex, name, x, y, w, h=0.54, fontsize=19):
    add_panel(slide, x, y, w, h, fill=rgb("f8fafc"))
    img = render_formula(latex, name, fontsize=fontsize, width=max(w, 3.0), height=max(h, 0.45))
    slide.shapes.add_picture(str(img), Inches(x + 0.12), Inches(y + 0.07), width=Inches(w - 0.24), height=Inches(h - 0.12))


def add_arrow(slide, x1, y1, x2, y2, color=LINE):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True


def load_metrics():
    data = json.loads(METRICS.read_text(encoding="utf-8"))
    summary = data["summary"]
    report = data["classification_report"]
    weakest = min(
        ((v["f1-score"], k, v["precision"], v["recall"], v["support"]) for k, v in report.items() if isinstance(v, dict) and "f1-score" in v),
        key=lambda x: x[0],
    )
    return summary, weakest


def make_deck():
    summary, weakest = load_metrics()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    def new_slide():
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        return s

    # 1
    slide = new_slide()
    add_title(slide, "RSCD视觉路面摩擦可供性估计：当前单模型算法解释", "基于刚才PPT扩展：架构、公式、变量、模块功能、代码实现与大白话说明", "单模型 / 全量评估")
    metrics = [
        ("完整测试 Top-1", f"{summary['top1']*100:.2f}%", TEAL),
        ("完整测试 Macro-F1", f"{summary['macro_f1']*100:.2f}%", BLUE),
        ("Mean-P / Mean-R", f"{summary['mean_precision']*100:.2f}% / {summary['mean_recall']*100:.2f}%", GREEN),
        (f"最弱类：{weakest[1]}", f"{weakest[0]*100:.2f}%", RED),
        ("参数量", "32.49M", AMBER),
    ]
    for i, (t, v, c) in enumerate(metrics):
        add_metric(slide, 0.55 + i * 2.50, 1.34, 2.22, t, v, c)
    add_bullet_panel(
        slide,
        0.72,
        2.42,
        5.95,
        3.80,
        "当前最好单模型是什么",
        [
            "ConvNeXt-Tiny 作为视觉主干，但在主干内部加入物理机制条件化的张量耦合。",
            "PhysicsTexture 显式提取湿滑、反光、暗水、雪/冰、纹理粗糙等摩擦线索。",
            "LocalPhysicsField 做弱分割式局部证据场，避免只靠整图平均。",
            "C3 token 负责 friction/material/roughness 的解耦与再耦合表达。",
            "hard-pair error gate 与 SRBR 只在诊断出的困难边界上做小范围校准。",
        ],
        TEAL,
        13.5,
    )
    add_bullet_panel(
        slide,
        6.95,
        2.42,
        5.55,
        3.80,
        "一句大白话",
        [
            "不是让 ConvNeXt 自己盲猜 27 类，而是先让它看懂整张路面，再把“湿不湿、亮不亮、粗不粗、是否被水膜抹平”等证据显式喂给模型。",
            "最终分类仍保持单模型；没有模型集成，也没有额外数据。",
            "当前指标来自全量测试集评估；该 checkpoint 不是完整训练集长训收敛版，formal full-manifest 训练用于验证更公平上限。",
        ],
        BLUE,
        13.5,
    )

    # 2
    slide = new_slide()
    add_title(slide, "任务与标签：27类其实是三类因素的组合", "先理解标签结构，后面的解耦/耦合设计才有意义", "RSCD-27")
    add_formula(slide, r"y=(f,m,r),\quad f\in\mathcal{F},\;m\in\mathcal{M},\;r\in\mathcal{R}", "label_factor", 0.78, 1.18, 5.85, 0.72, 21)
    add_formula(slide, r"\mathcal{Y}=\{(f,m,r)\;|\;(f,m,r)\ \mathrm{is\ a\ valid\ RSCD\ class}\}", "label_set", 6.82, 1.18, 5.60, 0.72, 19)
    add_panel_text(
        slide,
        0.78,
        2.08,
        3.75,
        3.95,
        "变量含义",
        "f：摩擦/路面状态，如 dry、wet、water、snow/ice 等。\n"
        "m：材质或非铺装状态，如 asphalt、concrete、gravel、mud。\n"
        "r：粗糙/严重程度，如 smooth、slight、severe；部分雪/泥类没有完整三因子。\n"
        "y：27类最终标签，本质是一个组合标签。",
        TEAL,
        13.2,
    )
    add_panel_text(
        slide,
        4.75,
        2.08,
        3.75,
        3.95,
        "为什么难",
        "同一个 roughness=slight 在 dry concrete、water concrete、wet asphalt 中视觉机制不同。\n"
        "水膜会遮挡细纹理；标线/强反光会制造伪粗糙；patch 视角不固定，不能假设下半部分一定是轮胎接触区。\n"
        "所以模型必须同时学单因素、两因素交互和三因素耦合。",
        BLUE,
        13.2,
    )
    add_panel_text(
        slide,
        8.72,
        2.08,
        3.70,
        3.95,
        "当前主要短板",
        "最弱类：water_concrete_slight，F1=76.04%。\n"
        "全量测试错误中，铺装类粗糙度相关错误占比最高，约 57.51%。\n"
        "Top混淆：dry_concrete_slight↔severe、water/wet_concrete smooth 边界、water_concrete_slight↔severe。",
        RED,
        13.2,
    )
    add_text(slide, 0.85, 6.35, 11.3, 0.52, "大白话：27类不是27个互不相关的名字，而是“路面状态 × 材质 × 粗糙程度”的组合；真正难点是组合后会出现新的视觉现象。", size=15, color=NAVY, bold=True)

    # 3
    slide = new_slide()
    add_title(slide, "整体架构：从一张 patch 到 27 类概率", "当前最好单模型的真实数据流：主干为锚，物理证据辅助，困难边界局部修正", "完整流程")
    boxes = [
        (0.55, 1.32, 1.85, "输入 x", "192×192\nImageNet归一化", TEAL),
        (2.75, 1.32, 2.05, "ConvNeXt Bθ", "768维全局视觉特征\n含早期张量耦合", BLUE),
        (5.15, 1.32, 2.05, "LocalField Lη", "8个局部软证据场\n0.08比例修正", GREEN),
        (7.55, 1.32, 2.20, "Physics Pφ/Aψ", "96维物理纹理\n64维语义物理注意", AMBER),
        (10.10, 1.32, 2.05, "Head + SRBR", "hard-pair门控\n输出27类logits", RED),
    ]
    for x, y, w, h1, h2, c in boxes:
        add_panel(slide, x, y, w, 0.88)
        add_text(slide, x + 0.08, y + 0.10, w - 0.16, 0.25, h1, size=12.5, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.08, y + 0.39, w - 0.16, 0.34, h2, size=10.2, color=SLATE, align=PP_ALIGN.CENTER)
    for x in [2.42, 4.88, 7.27, 9.90]:
        add_arrow(slide, x, 1.75, x + 0.25, 1.75)
    add_formula(slide, r"b=B_{\theta}(x),\qquad b'=b+\alpha A_{\eta}(L_{\eta}(x)),\quad \alpha=0.08", "overall_b", 0.75, 2.62, 5.75, 0.64, 18)
    add_formula(slide, r"z=\mathrm{Dropout}\!\left(\mathrm{LN}\left([\,b',P_{\phi}(x),A_{\psi}(x)\,]\right)\right)", "overall_z", 6.75, 2.62, 5.55, 0.64, 18)
    add_formula(slide, r"o=Wz+b_0+\Delta_{\mathrm{hp}}+\Delta_{\mathrm{srbr}},\qquad \hat{y}=\arg\max_c\ \mathrm{softmax}(o)_c", "overall_out", 1.50, 3.50, 10.25, 0.70, 18)
    add_panel_text(
        slide,
        0.75,
        4.55,
        5.55,
        1.65,
        "变量与模块",
        "Bθ：ConvNeXt主干；Lη：LocalPhysicsField；Aη：64→768适配器；Pφ：PhysicsTexture；Aψ：SemanticPhysicsAttention；LN：LayerNorm；W,b0：27类线性分类权重与偏置；Δhp/Δsrbr：困难边界的小范围特征/概率校准。",
        TEAL,
        12.2,
    )
    add_panel_text(
        slide,
        6.55,
        4.55,
        5.75,
        1.65,
        "程序实现",
        "主流程在 C3FaRNetSurfaceClassifier.forward：先 backbone(image)，再 local_physics_field_branch(image)，再 physics_branch / semantic_physics_attention_branch，拼接归一化后进入 linear head、hard-pair gate 与 SRBR。",
        BLUE,
        12.2,
    )
    add_text(slide, 0.82, 6.45, 11.3, 0.42, "大白话：先看整张路，再补充“湿、亮、暗水、纹理被抹平、粗糙颗粒”这些人能解释的线索；最后只在容易混淆的相邻类别边界上轻轻修正。", size=14.5, color=NAVY, bold=True)

    # 4
    slide = new_slide()
    add_title(slide, "PhysicsTexture：把摩擦相关的光学/纹理证据显式算出来", "这个模块是目前最稳定、最可解释的有效模块", "Pφ(x)")
    add_formula(slide, r"g=0.299R+0.587G+0.114B,\qquad s=\frac{\max(R,G,B)-\min(R,G,B)}{\max(R,G,B)+\varepsilon}", "phys_gray_sat", 0.70, 1.05, 11.75, 0.62, 17)
    add_formula(slide, r"\|\nabla g\|_2=\sqrt{(K_x*g)^2+(K_y*g)^2+\varepsilon},\qquad \Delta g=|K_{\Delta}*g|", "phys_grad_lap", 0.70, 1.82, 11.75, 0.62, 17)
    add_formula(slide, r"E_{\mathrm{wet}}=\mathrm{clip}(E_{\mathrm{specular}}+0.5E_{\mathrm{dark}},0,1),\qquad E_{\mathrm{rough}}=\sigma((\|\nabla g\|_2+0.5|\Delta g|-\tau)\kappa)", "phys_wet_rough", 0.70, 2.59, 11.75, 0.66, 16)
    add_panel_text(
        slide,
        0.75,
        3.55,
        3.75,
        2.12,
        "每个变量",
        "R,G,B：反归一化后的颜色通道。\ng：灰度/亮度图；s：饱和度。\nKx,Ky：Sobel一阶梯度核。\nKΔ：Laplacian二阶变化核。\nσ：sigmoid，把证据压到0到1。",
        TEAL,
        12.4,
    )
    add_panel_text(
        slide,
        4.75,
        3.55,
        3.75,
        2.12,
        "物理含义",
        "E_specular：亮且低饱和的镜面反光。\nE_dark：暗、平滑、低纹理的积水/水膜。\nE_wet：湿滑/水膜综合证据。\nE_rough：可见粗糙纹理或碎裂高频响应。",
        BLUE,
        12.4,
    )
    add_panel_text(
        slide,
        8.75,
        3.55,
        3.70,
        2.12,
        "程序怎么做",
        "src/friction_affordance/models/texture.py::PhysicsTextureBranch\n先算18个基础统计；quality_cues=True时扩展到42个统计；再经 LayerNorm→Linear→GELU→Linear 投影成96维。",
        AMBER,
        12.1,
    )
    add_text(slide, 0.82, 6.08, 11.4, 0.62, "大白话：它像一个“路面物理证据计量表”，不直接替代神经网络分类，而是把湿滑、反光、暗水、雪白、粗糙这些线索量化后交给后面的分类器。", size=14.5, color=NAVY, bold=True)

    # 5
    slide = new_slide()
    add_title(slide, "LocalPhysicsField：弱分割式局部物理场", "不是预测真实mask，而是保留局部证据的空间分布", "Lη(x)")
    add_formula(slide, r"\mathcal{F}(x)=\{F_1,\ldots,F_8\}=\{E_{\mathrm{spec}},E_{\mathrm{dark}},E_{\mathrm{film}},E_{\mathrm{erase}},E_{\mathrm{lowtex}},E_{\mathrm{rough}},E_{\mathrm{lowcon}},E_{\mathrm{mark}}\}", "local_fields", 0.62, 1.05, 12.0, 0.70, 15)
    add_formula(slide, r"u_j=\left[\mathrm{Pool}_{3\times3}(F_j),\ \mu(F_j),\ \sigma(F_j),\ \max(F_j),\ C(F_j)\right]", "local_stats", 0.62, 1.85, 12.0, 0.64, 17)
    add_formula(slide, r"C(F)=\frac{1}{HW}\sum_{i,j}F_{ij}\,[\mathrm{AvgPool}_{9\times9}(F)]_{ij},\qquad L_{\eta}(x)=\mathrm{MLP}(\mathrm{LN}([u_1,\ldots,u_8]))", "local_conn", 0.62, 2.62, 12.0, 0.68, 16)
    add_panel_text(
        slide,
        0.75,
        3.55,
        3.75,
        2.18,
        "变量含义",
        "Fj：第j个软证据场，数值越高代表该像素越像某种物理现象。\nPool3×3：把patch切成3×3格子做平均，保留“证据在哪里”。\nC(F)：软连通性，大片连续水膜会比零散噪点更高。",
        TEAL,
        12.2,
    )
    add_panel_text(
        slide,
        4.75,
        3.55,
        3.75,
        2.18,
        "实现逻辑",
        "LocalPhysicsFieldBranch 生成8个场；每个场产生9个格子均值+mean/std/max/connectedness。\n总统计维度为 8×(9+4)=104，经 LN→Linear→GELU→Linear 得到64维。",
        BLUE,
        12.2,
    )
    add_panel_text(
        slide,
        8.75,
        3.55,
        3.70,
        2.18,
        "为什么不用底部ROI",
        "RSCD多是patch图，拍摄方向不总是正前方，转弯或局部裁剪会破坏“下方=接触区”的假设。\n因此当前只保留网格分布，不强行绑定底部区域。",
        AMBER,
        12.2,
    )
    add_text(slide, 0.82, 6.10, 11.4, 0.62, "大白话：PhysicsTexture看“整张图总体像不像湿滑/粗糙”，LocalPhysicsField看“这些证据是成片出现、局部出现，还是可能只是标线/噪声”。", size=14.5, color=NAVY, bold=True)

    # 6
    slide = new_slide()
    add_title(slide, "早期机制条件化主干：把物理机制放进ConvNeXt内部", "关键判断：后期head/residual太晚，真正有潜力的是早期/中期特征机制化", "Bθ(x)")
    add_formula(slide, r"X_{\ell+1}=C_{\ell}(X_{\ell})+\Delta_{\ell}", "backbone_update", 0.80, 1.08, 5.55, 0.62, 21)
    add_formula(slide, r"\Delta_{\ell}=\Pi_{\Vert\cdot\Vert\leq r}\left[s\cdot g_{\ell}(x)\cdot\mathcal{T}_{\ell}(X_{\ell};e)\right]", "backbone_delta", 6.65, 1.08, 5.80, 0.62, 18)
    add_formula(slide, r"g_{\ell}(x)=\sigma(h_{\ell}(X_{\ell},e))\cdot q_{\mathrm{phys}}(e)\cdot q_{\mathrm{artifact}}(e)", "backbone_gate", 1.55, 1.92, 10.05, 0.62, 18)
    add_panel_text(
        slide,
        0.75,
        2.90,
        3.65,
        2.42,
        "变量解释",
        "Xℓ：第ℓ层ConvNeXt特征图。\nCℓ：原始ConvNeXt层。\nΔℓ：物理机制张量耦合带来的微小特征更新。\nΠ：范数裁剪，限制更新幅度。\ne：16维低层物理证据。",
        TEAL,
        12.2,
    )
    add_panel_text(
        slide,
        4.60,
        2.90,
        3.85,
        2.42,
        "模块功能",
        "只让与RSCD困难组合相关的cell被激活，例如 wet/water concrete 的 film-roughness 单元。\n更新发生在主干特征图阶段，而不是最终logits之后。\n零初始化/小尺度保证从强锚点平滑开始。",
        BLUE,
        12.2,
    )
    add_panel_text(
        slide,
        8.65,
        2.90,
        3.70,
        2.42,
        "程序实现",
        "src/friction_affordance/models/backbone.py\nConvNeXtGateCalibratedTensorCouplingBackbone 在指定 features 层后插入 _GateCalibratedTensorCouplingBank；当前配置用 concrete_film_rough_stem 变体。",
        AMBER,
        12.0,
    )
    add_text(slide, 0.82, 5.72, 11.35, 0.72, "大白话：不是等模型已经做完判断后再改答案，而是在它提取纹理和边缘时就提醒它：这个地方可能是水膜遮住了混凝土粗糙度，别把“看起来平滑”简单当成真正smooth。", size=14.2, color=NAVY, bold=True)

    # 7
    slide = new_slide()
    add_title(slide, "C3解耦/再耦合：让模型理解组合标签", "把摩擦、材质、粗糙度拆开学，再允许它们形成组合后的新现象", "C3 tokens")
    add_formula(slide, r"(t_f,t_m,t_r,t_c)=D_{\omega}(z,e)", "c3_tokens", 0.70, 1.02, 5.40, 0.62, 21)
    add_formula(slide, r"o^{\mathrm{C3}}_{fmr}=A_f+B_m+C_r+D_{fm}+E_{fr}+G_{mr}+H_{fmr}", "c3_score", 6.35, 1.02, 6.05, 0.62, 18)
    add_formula(slide, r"j=[\,t_c,\ e,\ \rho_R\,],\qquad \rho_R=\sigma(r_{\xi}(e,t_c))", "c3_judge", 1.15, 1.82, 10.95, 0.62, 18)
    add_panel_text(
        slide,
        0.75,
        2.85,
        3.75,
        2.55,
        "变量解释",
        "tf/tm/tr：摩擦、材质、粗糙度token。\ntc：耦合token，专门表示组合后现象。\nAf/Bm/Cr：一阶因素项。\nDfm/Efr/Gmr：两因素交互项。\nHfmr：三因素耦合项。",
        TEAL,
        12.1,
    )
    add_panel_text(
        slide,
        4.75,
        2.85,
        3.75,
        2.55,
        "当前模型里的作用",
        "当前最好head是 hardpair_error_gated_calibrated，主logits由线性头提供。\nC3 token、ρR和e主要组成 judge_input，用来判断困难相邻类别边界是否需要校准。\n这比直接把C3 logits硬加到结果上更稳。",
        BLUE,
        12.1,
    )
    add_panel_text(
        slide,
        8.75,
        2.85,
        3.70,
        2.55,
        "程序实现",
        "FactorQueryDecoder：把 z 和 e 解码为4个token。\nRoughnessReliability：估计可见粗糙度是否可信。\nCoupledTensorHead：实现一阶、二阶和三阶耦合张量打分。",
        AMBER,
        12.1,
    )
    add_text(slide, 0.82, 5.78, 11.35, 0.62, "大白话：模型先分别问“它湿吗、是什么材质、粗不粗”，再问“湿+混凝土+轻微粗糙放在一起会不会因为水膜而长得完全不一样”。", size=14.5, color=NAVY, bold=True)

    # 8
    slide = new_slide()
    add_title(slide, "困难边界校准：hard-pair error gate 与 SRBR", "只在相邻、近边界、物理证据匹配时动手，避免救弱类时伤强类", "边界机制")
    add_formula(slide, r"g_{ij}=\sigma((m-|o_i-o_j|)\tau)\,(p_i+p_j)\,g_{\mathrm{sample}}\,g_{\mathrm{err}}\,q_{\mathrm{phys}}", "hp_gate", 0.62, 1.02, 12.0, 0.62, 16)
    add_formula(slide, r"\delta_{ij}=s_{ij}\lambda_{\mathrm{hp}}\,g_{ij}\tanh(h_{ij}(j)),\qquad o_i\leftarrow o_i+\delta_{ij},\quad o_j\leftarrow o_j-\delta_{ij}", "hp_delta", 0.62, 1.78, 12.0, 0.62, 16)
    add_formula(slide, r"\Delta z_{s\to t}=\gamma\,q_{\mathrm{region}}\,q_{\mathrm{close}}\,q_{\mathrm{phys}}\,\rho_s\,r_{\chi}(c)\,\frac{w_t-w_s}{\|w_t-w_s\|_2}", "srbr_delta", 0.62, 2.54, 12.0, 0.62, 15)
    add_panel_text(
        slide,
        0.75,
        3.48,
        3.75,
        2.25,
        "hard-pair变量",
        "i,j：一对容易混淆的相邻类。\n|oi-oj|：两类logit差，越小越接近边界。\npi+pj：样本落在这对类别附近的概率质量。\ngerr：学习到的锚点错误门。\nqphys：物理证据门。",
        TEAL,
        12.0,
    )
    add_panel_text(
        slide,
        4.75,
        3.48,
        3.75,
        2.25,
        "SRBR变量",
        "s/t：source/target类别。\nqregion：source被预测且target进入Top-k。\nqclose：source-target margin足够近。\nρs：source类验证可靠度。\nwt-ws：分类器判别方向，不自由改logits。",
        BLUE,
        12.0,
    )
    add_panel_text(
        slide,
        8.75,
        3.48,
        3.70,
        2.25,
        "当前有效路线",
        "SRBR目前保留的有效路线：dry_concrete_smooth → dry_concrete_slight。\n配置：topk=3、margin=1.0、source_f1=0.9073、physics_gate_floor=0.5、scale=7.0。",
        AMBER,
        12.0,
    )
    add_text(slide, 0.82, 6.05, 11.35, 0.62, "大白话：模型只有在“它本来就犹豫，而且物理证据也支持这个边界修正”时才轻推一下；不是把所有样本都暴力改答案。", size=14.5, color=NAVY, bold=True)

    # 9
    slide = new_slide()
    add_title(slide, "训练目标：修弱类，同时保护强锚点", "当前正式训练配置采用锚点一致性与非退化保护，避免局部提升换来整体下降", "Loss")
    add_formula(slide, r"\mathcal{L}=\mathcal{L}_{\mathrm{CE}}+\lambda_{\mathrm{KD}}T^2D_{\mathrm{KL}}(p_T^{(T)}\parallel p_S^{(T)})+\lambda_{\mathrm{NF}}\mathcal{L}_{\mathrm{no\!-\!flip}}+\lambda_g\mathcal{L}_{\mathrm{gate}}+\mathcal{L}_{\mathrm{focus}}", "loss_total", 0.58, 1.05, 12.15, 0.70, 14)
    add_formula(slide, r"\mathcal{L}_{\mathrm{CE}}=-\log p_S(y),\qquad p_S=\mathrm{softmax}(o_S),\qquad p_T^{(T)}=\mathrm{softmax}(o_T/T)", "loss_ce", 0.92, 1.92, 11.45, 0.62, 17)
    add_panel_text(
        slide,
        0.75,
        2.92,
        3.75,
        2.60,
        "每一项的含义",
        "LCE：27类主监督。\nKD/KL：学生模型不要偏离强教师/锚点太远。\nLno-flip：教师高置信且正确的非重点类不要被翻错。\nLgate：让错误门学会在哪里开。\nLfocus：弱类额外加权。",
        TEAL,
        12.2,
    )
    add_panel_text(
        slide,
        4.75,
        2.92,
        3.75,
        2.60,
        "当前关键权重",
        "anchor_consistency_weight=0.55。\nanchor_no_flip_weight=0.22。\nanchor_error_gate_weight=0.10。\nfocus_ce_extra_weight=0.30。\n蒸馏温度 T=2.0。\n重点类：water/wet concrete slight/severe。",
        BLUE,
        12.2,
    )
    add_panel_text(
        slide,
        8.75,
        2.92,
        3.70,
        2.60,
        "程序怎么训练",
        "train.py → c3_experiment.py。\nAMP开启；batch_size=8，grad_accum=2。\n当前只放开少量前缀训练：gate_calibrated_tensor_coupling_banks、pairwise_hardpair_experts、pairwise_hardpair_error_gates。",
        AMBER,
        12.1,
    )
    add_text(slide, 0.82, 5.95, 11.35, 0.72, "大白话：训练时一边告诉模型“这些弱类要重点学”，一边拉住它“别把原来已经会的强类弄坏”。这就是为什么不能只追求某个弱类F1，而要看Macro-F1和no-harm。", size=14.2, color=NAVY, bold=True)

    # 10
    slide = new_slide()
    add_title(slide, "当前结论与下一步：哪些模块真实有效", "这页用于汇报收束：指标、失败路线、下一步最有希望的方向", "总结")
    add_bullet_panel(
        slide,
        0.75,
        1.18,
        3.65,
        2.15,
        "当前最好公平指标",
        [
            f"Top-1：{summary['top1']*100:.2f}%，Macro-F1：{summary['macro_f1']*100:.2f}%。",
            f"Mean-P/Mean-R：{summary['mean_precision']*100:.2f}% / {summary['mean_recall']*100:.2f}%。",
            f"最弱类：{weakest[1]}，F1={weakest[0]*100:.2f}%。",
            "参数量：32.49M；测试集 N=49,500。",
        ],
        TEAL,
        11.8,
    )
    add_bullet_panel(
        slide,
        4.60,
        1.18,
        3.85,
        2.15,
        "目前最有效部分",
        [
            "PhysicsTexture：提供稳定、可解释的湿滑/粗糙证据。",
            "LocalPhysicsField：补充局部分布，减少整图平均损失。",
            "早期张量耦合主干：比后期head更符合瓶颈。",
            "SRBR：对已验证可靠边界做有保护的特征路由。",
        ],
        BLUE,
        11.8,
    )
    add_bullet_panel(
        slide,
        8.65,
        1.18,
        3.75,
        2.15,
        "失败/谨慎路线",
        [
            "VLM语义描述：细粒度粗糙边界不稳定，成本高。",
            "晚期residual/head/aux：常救少数类但伤其他类。",
            "固定底部ROI：RSCD patch视角不固定，容易引入偏置。",
            "泛化loss surgery：瓶颈更像特征机制，不只是梯度冲突。",
        ],
        RED,
        11.8,
    )
    add_bullet_panel(
        slide,
        0.75,
        3.75,
        5.65,
        2.35,
        "下一步最有希望",
        [
            "等待formal full-manifest单模型完整训练结果，确认当前结构在全训练清单上的公平上限。",
            "继续把有效物理机制前移到早期/中期主干，而不是在最终logits后补丁。",
            "重点攻克 water/wet concrete slight/severe：显式建模可见粗糙、隐藏粗糙、水膜遮挡、伪粗糙。",
            "每个新机制都必须做RSCD任务适配与同预算no-harm对照。",
        ],
        TEAL,
        12.3,
    )
    add_bullet_panel(
        slide,
        6.70,
        3.75,
        5.70,
        2.35,
        "汇报时可以这样说",
        [
            "当前路线已经从“堆模块”转向“根据RSCD三因子耦合机理设计模块”。",
            "指标短板不是整体不会分类，而是水膜/混凝土/粗糙度耦合边界仍不够精确。",
            "下一步不再盲目加head，而是围绕边界证据可靠性做早期机制条件化。",
        ],
        BLUE,
        12.3,
    )
    add_text(slide, 0.82, 6.38, 11.35, 0.38, "大白话：现在最该做的是把模型“看纹理的方式”改对，而不是在它已经猜完以后再贴补丁。", size=15, color=NAVY, bold=True)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = make_deck()
    print(out)
